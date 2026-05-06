from pptx import Presentation
from pptx.util import Inches, Pt
import os

prs = Presentation()

# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Cache-Aware Graph Reordering for Efficient GNN Training"
slide.placeholders[1].text = "Sahil Murtaza, Youke Zhang, Zach Yu, Albert Liu\nUniversity of Virginia\nCS 6501 · Graph ML"

# Slide 2: Motivation
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Motivation: GNNs Hit a Memory Wall on GPUs"
tf = slide.placeholders[1].text_frame
tf.text = "The Bottleneck: GPUs optimized for dense workloads; GNN scatter-gather defeats HBM coalescing."
tf.add_paragraph().text = "The Scale Problem: Large feature matrices (e.g., 960MB) exceed L2 cache capacity (40MB)."
tf.add_paragraph().text = "Root Cause: Random node ordering causes cache misses on neighbor fetches."

# Slide 3: Core Idea
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Core Idea: Node Reordering = Lossless Locality Boost"
tf = slide.placeholders[1].text_frame
tf.text = "Key Insight: GNN message-passing is equivariant to node relabeling (accuracy guaranteed identical)."
tf.add_paragraph().text = "Bandwidth Minimization: Permutation minimizes matrix bandwidth, clustering non-zeros."
tf.add_paragraph().text = "Pre-processing: One-time O(|E| log|E|) operation; no changes to training pipeline."

# Slide 4: Baseline Methods
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Baseline Methods: RCM & METIS"
tf = slide.placeholders[1].text_frame
tf.text = "Reverse Cuthill-McKee (RCM): Produces a band-diagonal adjacency matrix."
tf.add_paragraph().text = "METIS Graph Partitioning: Produces a block-diagonal adjacency matrix."
tf.add_paragraph().text = "Prior work applies these independently with blind hyperparameter sweeps."

# Slide 5: Novel Contribution 1
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Novel Contribution 1: Two-Level Hybrid Reordering"
tf = slide.placeholders[1].text_frame
tf.text = "Stage 1: METIS Partitioning (groups nodes by community)."
tf.add_paragraph().text = "Stage 2: Intra-block RCM (applies RCM within each partition)."
tf.add_paragraph().text = "Result: Block-diagonal + band-diagonal structure for multiplicative cache win."

# Slide 6: Novel Contribution 2
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Novel Contribution 2: Cache-Size-Aware Partition Selection"
tf = slide.placeholders[1].text_frame
tf.text = "Insight: Feature sub-matrix must fit in L2 cache."
tf.add_paragraph().text = "Derived Formula: k* = ceil(n * d * b / C)"
tf.add_paragraph().text = "Eliminates blind grid search; automatically adapts to any GPU architecture."

# Slide 7: Experimental Setup
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Experimental Setup: Datasets, Models, and Hardware"
tf = slide.placeholders[1].text_frame
tf.text = "OGB Datasets: ogbn-arxiv (47% L2 coverage), ogbn-products (4% L2 coverage)."
tf.add_paragraph().text = "GNN Models: GraphSAGE, GAT (2x neighbor feature accesses)."
tf.add_paragraph().text = "Hardware: 2x NVIDIA A100-SXM4-80GB (40 MB L2 cache per GPU)."

# Slide 8: Progress Report
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Progress Report: Baseline Measurements & Proxy Metrics"
tf = slide.placeholders[1].text_frame
tf.text = "GAT backward pass dominates latency due to per-edge attention backprop."
tf.add_paragraph().text = "Proxy Metrics used due to hardware counter restrictions:"
tf.add_paragraph().text = "1. Analytical Cache Coverage (ACC)"
tf.add_paragraph().text = "2. Temporal Reuse Ratio (TRR)"

# Slide 9: Project Roadmap
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Project Roadmap: Current Status & Next Steps"
tf = slide.placeholders[1].text_frame
tf.text = "Phase 1: Environment & Baselines (Done)"
tf.add_paragraph().text = "Phase 2: Baseline Reordering (In Progress)"
tf.add_paragraph().text = "Phase 3: Novel Methods Implementation (To Do)"
tf.add_paragraph().text = "Phase 4: Full Profiling & Ablation (To Do)"
tf.add_paragraph().text = "Target: 15-30% latency reduction."

# Slide 10: Summary
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Summary"
tf = slide.placeholders[1].text_frame
tf.text = "Problem: GNN training bottlenecked by poor L2 cache locality."
tf.add_paragraph().text = "Approach: Lossless node reordering to cluster connected nodes."
tf.add_paragraph().text = "Contributions: Hybrid reordering (METIS + RCM) and cache-aware k* selection."

output_path = "/mnt/data/Cache_Aware_Graph_Reordering.pptx"
prs.save(output_path)
print(output_path)